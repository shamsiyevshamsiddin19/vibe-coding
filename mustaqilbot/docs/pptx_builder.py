"""JSON slayd ma'lumotidan professional PPTX fayl yaratadi.

Dizayn dvigateli:
- Gradient: diagonal / vertikal / 3-rangli / radial — har slaydda navbatlashadi
- Bezaklar: shaffof doiralar, halqalar, uchburchak, lenta, nuqtalar to'ri
- 6 xil rang-mavzu (mavzu nomidan tanlanadi)
- 12 xil maket: cover, content, two_col, stat, quote, chart, image, end,
  section, steps, cards, timeline

Pexels orqali mavzuga mos fon-foto (ixtiyoriy — kalit/rasm topilmasa gradient).
Kesilgan/buzuq JSON'ga chidamli; 3 dan kam slayd chiqsa XATO ko'tariladi.
"""
from __future__ import annotations
import asyncio
import io
import os
import tempfile
from datetime import date

from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn

from docs.json_utils import parse_json_objects

_SLIDE_W = Inches(13.33)
_SLIDE_H = Inches(7.5)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_OVERLAY = RGBColor(0x0a, 0x0a, 0x14)
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

_C = RGBColor.from_string

# Rang-mavzular: bg (eng to'q) → bg2/bg3 (oraliq) → accent/accent2 (yorqin),
# ring (xira aksent), title/text, card/card2 (kartochka fonlari)
_THEMES = [
    {  # 1. Ko'k-professional
        "bg": _C("0B1220"), "bg2": _C("15213B"), "bg3": _C("0E2A4A"),
        "accent": _C("38BDF8"), "accent2": _C("60A5FA"), "ring": _C("1E3A5F"),
        "title": _WHITE, "text": _C("CBD5E1"), "card": _C("16223C"), "card2": _C("1E2E4E"),
    },
    {  # 2. Yashil-akademik
        "bg": _C("07160F"), "bg2": _C("0D2A1F"), "bg3": _C("0A3325"),
        "accent": _C("34D399"), "accent2": _C("6EE7B7"), "ring": _C("155040"),
        "title": _WHITE, "text": _C("D1FAE5"), "card": _C("0F3226"), "card2": _C("15402F"),
    },
    {  # 3. Binafsha-zamonaviy
        "bg": _C("140F24"), "bg2": _C("231A3C"), "bg3": _C("2A1E4A"),
        "accent": _C("C084FC"), "accent2": _C("A78BFA"), "ring": _C("4C2A7A"),
        "title": _WHITE, "text": _C("E9D5FF"), "card": _C("241A40"), "card2": _C("2E2250"),
    },
    {  # 4. Marjon-oltin (issiq)
        "bg": _C("1A1210"), "bg2": _C("2A1B16"), "bg3": _C("33201A"),
        "accent": _C("FB923C"), "accent2": _C("FBBF24"), "ring": _C("7C3A18"),
        "title": _WHITE, "text": _C("FCE7D8"), "card": _C("2B1E18"), "card2": _C("38271F"),
    },
    {  # 5. Feruza-cyan
        "bg": _C("07171A"), "bg2": _C("0C2B30"), "bg3": _C("0A343B"),
        "accent": _C("22D3EE"), "accent2": _C("2DD4BF"), "ring": _C("155E67"),
        "title": _WHITE, "text": _C("CFFAFE"), "card": _C("0F3238"), "card2": _C("154047"),
    },
    {  # 6. Qizil-oltin (jasur)
        "bg": _C("150A0A"), "bg2": _C("2A1414"), "bg3": _C("331717"),
        "accent": _C("F87171"), "accent2": _C("FBBF24"), "ring": _C("7A2323"),
        "title": _WHITE, "text": _C("FEE2E2"), "card": _C("261212"), "card2": _C("331A1A"),
    },
]


def _theme_for(topic: str) -> dict:
    return _THEMES[sum(ord(c) for c in (topic or "a")) % len(_THEMES)]


# ─────────────────────── Quyi-daraja yordamchilar ───────────────────────

def _set_alpha(shape, pct: int):
    try:
        solid_fill = shape.fill._xPr.find(qn("a:solidFill"))
        if solid_fill is None:
            return
        srgb = solid_fill.find(qn("a:srgbClr"))
        if srgb is None:
            return
        srgb.append(srgb.makeelement(qn("a:alpha"),
                                     {"val": str(max(0, min(100, pct)) * 1000)}))
    except Exception:
        pass


def _line_alpha(shape, pct: int):
    try:
        ln = shape.line._get_or_add_ln()
        srgb = ln.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        srgb.append(srgb.makeelement(qn("a:alpha"),
                                     {"val": str(max(0, min(100, pct)) * 1000)}))
    except Exception:
        pass


def _rect(slide, left, top, width, height, shape_type=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape_type, left, top, width, height)
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _fill(shape, color: RGBColor, alpha: int | None = None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if alpha is not None:
        _set_alpha(shape, alpha)


def _send_to_back(slide, shape):
    tree = slide.shapes._spTree
    tree.remove(shape._element)
    tree.insert(2, shape._element)


def _text(slide, text, left, top, width, height, size=18, bold=False,
         color=_WHITE, align=PP_ALIGN.LEFT, wrap=True, anchor=None):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = wrap
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tx


# ─────────────────────── Gradient dvigateli ───────────────────────

def _grad_xml(stops, angle=45.0, radial=False, center=(50, 50)) -> str:
    """stops: [(pos_0_100, RGBColor, alpha_0_100|None), ...]"""
    gsl = ""
    for pos, color, alpha in stops:
        a = f'<a:alpha val="{int(alpha * 1000)}"/>' if alpha is not None else ""
        gsl += (f'<a:gs pos="{int(pos * 1000)}">'
                f'<a:srgbClr val="{str(color)}">{a}</a:srgbClr></a:gs>')
    if radial:
        cx, cy = center
        l, t = int(cx * 1000), int(cy * 1000)
        r, b = int((100 - cx) * 1000), int((100 - cy) * 1000)
        path = f'<a:path path="circle"><a:fillToRect l="{l}" t="{t}" r="{r}" b="{b}"/></a:path>'
    else:
        path = f'<a:lin ang="{int(angle * 60000) % 21600000}" scaled="1"/>'
    return (f'<a:gradFill xmlns:a="{_A_NS}" rotWithShape="1">'
            f'<a:gsLst>{gsl}</a:gsLst>{path}</a:gradFill>')


def _apply_grad(shape, xml: str):
    spPr = shape._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill",
                "a:pattFill", "a:grpFill"):
        e = spPr.find(qn(tag))
        if e is not None:
            spPr.remove(e)
    grad = parse_xml(xml)
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        ln.addprevious(grad)
    else:
        spPr.append(grad)


# Har slaydda navbatlashadigan gradient uslublari
_GRAD_ROT = ["diag_tri", "vert", "radial_corner", "diag2", "radial_center", "vert_tri"]
_DECOR_ROT = ["blobs", "wedge", "rings", "none", "dots", "band"]


def _bg_gradient(slide, theme, style: str):
    rect = _rect(slide, 0, 0, _SLIDE_W, _SLIDE_H)
    bg, bg2, bg3 = theme["bg"], theme["bg2"], theme["bg3"]
    try:
        if style == "diag_tri":
            xml = _grad_xml([(0, bg, None), (55, bg3, None), (100, bg2, None)], angle=52)
        elif style == "vert":
            xml = _grad_xml([(0, bg2, None), (100, bg, None)], angle=90)
        elif style == "radial_corner":
            xml = _grad_xml([(0, bg3, None), (100, bg, None)], radial=True, center=(16, 14))
        elif style == "diag2":
            xml = _grad_xml([(0, bg, None), (100, bg3, None)], angle=125)
        elif style == "radial_center":
            xml = _grad_xml([(0, bg2, None), (100, bg, None)], radial=True, center=(50, 40))
        else:  # vert_tri
            xml = _grad_xml([(0, bg, None), (50, bg2, None), (100, bg, None)], angle=90)
        _apply_grad(rect, xml)
    except Exception:
        _fill(rect, bg)
    _send_to_back(slide, rect)
    return rect


def _decorate(slide, theme, style: str):
    accent, accent2, ring = theme["accent"], theme["accent2"], theme["ring"]
    try:
        if style == "blobs":
            c1 = _rect(slide, Inches(9.6), Inches(-1.3), Inches(5.2), Inches(5.2), MSO_SHAPE.OVAL)
            _fill(c1, accent2, alpha=9)
            c2 = _rect(slide, Inches(-1.6), Inches(4.6), Inches(4.4), Inches(4.4), MSO_SHAPE.OVAL)
            _fill(c2, accent, alpha=8)
        elif style == "wedge":
            w = _rect(slide, Inches(10.9), 0, Inches(2.9), Inches(2.9), MSO_SHAPE.RIGHT_TRIANGLE)
            w.rotation = 90
            _fill(w, accent, alpha=12)
        elif style == "rings":
            for rr in (5.8, 4.3, 2.8):
                ring_sp = _rect(slide, Inches(10.6 - rr / 2), Inches(6.2 - rr / 2),
                                Inches(rr), Inches(rr), MSO_SHAPE.OVAL)
                ring_sp.fill.background()
                ring_sp.line.color.rgb = ring
                ring_sp.line.width = Pt(1.1)
        elif style == "band":
            b = _rect(slide, Inches(-2.5), Inches(5.4), Inches(19), Inches(1.5))
            b.rotation = -7
            _fill(b, accent, alpha=7)
        elif style == "dots":
            for r in range(4):
                for c in range(5):
                    d = _rect(slide, Inches(11.15 + c * 0.4), Inches(0.55 + r * 0.4),
                              Inches(0.12), Inches(0.12), MSO_SHAPE.OVAL)
                    _fill(d, accent, alpha=24)
        # "none" → bezaksiz (nafas olish uchun)
    except Exception:
        pass


def _scene(slide, theme, idx: int):
    """Har slaydga navbatlashuvchi gradient + bezak (xilma-xillik)."""
    _bg_gradient(slide, theme, _GRAD_ROT[idx % len(_GRAD_ROT)])
    _decorate(slide, theme, _DECOR_ROT[(idx * 2 + 1) % len(_DECOR_ROT)])


def _cover_photo(slide, image_bytes: bytes, overlay_alpha: int = 58):
    im = PILImage.open(io.BytesIO(image_bytes))
    iw, ih = im.size
    box_ratio = _SLIDE_W / _SLIDE_H
    img_ratio = iw / ih
    pic = slide.shapes.add_picture(io.BytesIO(image_bytes), 0, 0,
                                   width=_SLIDE_W, height=_SLIDE_H)
    if img_ratio > box_ratio:
        crop = (1 - box_ratio / img_ratio) / 2
        pic.crop_left, pic.crop_right = crop, crop
    else:
        crop = (1 - img_ratio / box_ratio) / 2
        pic.crop_top, pic.crop_bottom = crop, crop
    overlay = _rect(slide, 0, 0, _SLIDE_W, _SLIDE_H)
    _fill(overlay, _OVERLAY, alpha=overlay_alpha)
    return pic


def _accent_bar(slide, theme, top, height=Inches(0.06)):
    bar = _rect(slide, 0, top, _SLIDE_W, height)
    _fill(bar, theme["accent"])
    return bar


def _page_number(slide, theme, idx, total):
    _text(slide, f"{idx} / {total}", Inches(12.0), Inches(7.02), Inches(1.1), Inches(0.4),
         size=11, color=theme["text"], align=PP_ALIGN.RIGHT)


def _title_block(slide, theme, title, top=0.85, size=29):
    _text(slide, title, Inches(0.6), Inches(top), Inches(12.1), Inches(1.0),
         size=size, bold=True, color=theme["title"])
    rect = _rect(slide, Inches(0.65), Inches(top + 0.82), Inches(1.5), Inches(0.045))
    _fill(rect, theme["accent"])


def _lines_for(text: str, chars_per_line: int) -> int:
    return max(1, (len(text) + chars_per_line - 1) // chars_per_line)


def _new(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ─────────────────────── Maketlar (layouts) ───────────────────────

def _slide_cover(prs, theme, sd, image_bytes, author, subject):
    slide = _new(prs)
    if image_bytes:
        _cover_photo(slide, image_bytes, overlay_alpha=58)
    else:
        _bg_gradient(slide, theme, "diag_tri")
        _decorate(slide, theme, "rings")
        _accent_bar(slide, theme, Inches(3.35))
    title = str(sd.get("title", ""))[:140]
    _text(slide, title, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.9),
         size=38, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
    if subject:
        _text(slide, subject, Inches(0.9), Inches(4.05), Inches(11.5), Inches(0.5),
             size=17, color=theme["accent"], align=PP_ALIGN.CENTER)
    _text(slide, author, Inches(0.9), Inches(4.6), Inches(11.5), Inches(0.5),
         size=18, color=RGBColor(0xe6, 0xe6, 0xeb), align=PP_ALIGN.CENTER)
    _text(slide, str(date.today().year), Inches(0.9), Inches(5.1), Inches(11.5), Inches(0.4),
         size=14, color=RGBColor(0xc8, 0xc8, 0xd2), align=PP_ALIGN.CENTER)


def _slide_content(prs, theme, sd, idx, total):
    slide = _new(prs)
    _scene(slide, theme, idx)
    _title_block(slide, theme, str(sd.get("title", ""))[:110])
    points = [str(p) for p in sd.get("points", []) if str(p).strip()][:5]
    top = 2.15
    for pt in points:
        h = 0.5 + 0.3 * (_lines_for(pt, 68) - 1)
        card = _rect(slide, Inches(0.6), Inches(top), Inches(12.1), Inches(h),
                     shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0] = 0.09
        _fill(card, theme["card"])
        strip = _rect(slide, Inches(0.6), Inches(top), Inches(0.09), Inches(h),
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        strip.adjustments[0] = 0.5
        _fill(strip, theme["accent"])
        _text(slide, pt, Inches(0.98), Inches(top), Inches(11.5), Inches(h),
             size=17, color=theme["text"], anchor=MSO_ANCHOR.MIDDLE)
        top += h + 0.18
        if top > 6.7:
            break
    _page_number(slide, theme, idx, total)


def _slide_section(prs, theme, sd, idx, total, section_no):
    """Bo'lim ajratuvchi: chapda rangli panel + katta raqam, o'ngda sarlavha."""
    slide = _new(prs)
    _bg_gradient(slide, theme, "diag2")
    _decorate(slide, theme, "dots")
    panel = _rect(slide, 0, 0, Inches(4.6), _SLIDE_H)
    _fill(panel, theme["accent"])
    _text(slide, f"{section_no:02d}", Inches(0.3), Inches(2.0), Inches(4.0), Inches(2.8),
         size=150, bold=True, color=theme["bg"], align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
    title = str(sd.get("title", ""))[:90]
    _text(slide, title, Inches(5.1), Inches(2.7), Inches(7.6), Inches(2.2),
         size=34, bold=True, color=theme["title"], anchor=MSO_ANCHOR.MIDDLE)
    sub = [str(p) for p in sd.get("points", []) if str(p).strip()][:1]
    if sub:
        _text(slide, sub[0], Inches(5.1), Inches(4.7), Inches(7.6), Inches(1.2),
             size=17, color=theme["text"])
    _page_number(slide, theme, idx, total)


def _slide_steps(prs, theme, sd, idx, total):
    """Raqamlangan bosqichlar (jarayon/algoritm) — gorizontal kartochkalar."""
    slide = _new(prs)
    _scene(slide, theme, idx)
    _title_block(slide, theme, str(sd.get("title", ""))[:110])
    steps = [str(p) for p in sd.get("points", []) if str(p).strip()][:4]
    n = max(1, len(steps))
    gap = 0.35
    total_w = 12.1
    cw = (total_w - gap * (n - 1)) / n
    left = 0.6
    for i, st in enumerate(steps):
        x = left + i * (cw + gap)
        card = _rect(slide, Inches(x), Inches(2.5), Inches(cw), Inches(3.7),
                     shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0] = 0.06
        _fill(card, theme["card"])
        # Raqam doirasi
        circ = _rect(slide, Inches(x + cw / 2 - 0.55), Inches(2.15), Inches(1.1), Inches(1.1),
                     MSO_SHAPE.OVAL)
        _fill(circ, theme["accent"])
        _text(slide, str(i + 1), Inches(x + cw / 2 - 0.55), Inches(2.15), Inches(1.1), Inches(1.1),
             size=32, bold=True, color=theme["bg"], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, st, Inches(x + 0.2), Inches(3.5), Inches(cw - 0.4), Inches(2.5),
             size=15, color=theme["text"], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    _page_number(slide, theme, idx, total)


def _slide_cards(prs, theme, sd, idx, total):
    """Kartochka to'ri (2×2 yoki 2×N) — kategoriyalar/xususiyatlar."""
    slide = _new(prs)
    _scene(slide, theme, idx)
    _title_block(slide, theme, str(sd.get("title", ""))[:110])
    pts = [str(p) for p in sd.get("points", []) if str(p).strip()][:4]
    cols = 2
    gap = 0.4
    cw = (12.1 - gap) / cols
    ch = 1.95
    for i, pt in enumerate(pts):
        r, c = divmod(i, cols)
        x = 0.6 + c * (cw + gap)
        y = 2.25 + r * (ch + 0.35)
        card = _rect(slide, Inches(x), Inches(y), Inches(cw), Inches(ch),
                     shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0] = 0.07
        _fill(card, theme["card"])
        num = _rect(slide, Inches(x + 0.3), Inches(y + 0.3), Inches(0.7), Inches(0.7),
                    MSO_SHAPE.ROUNDED_RECTANGLE)
        num.adjustments[0] = 0.3
        _fill(num, theme["accent"])
        _text(slide, str(i + 1), Inches(x + 0.3), Inches(y + 0.3), Inches(0.7), Inches(0.7),
             size=22, bold=True, color=theme["bg"], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, pt, Inches(x + 1.2), Inches(y + 0.2), Inches(cw - 1.5), Inches(ch - 0.4),
             size=15, color=theme["text"], anchor=MSO_ANCHOR.MIDDLE)
    _page_number(slide, theme, idx, total)


def _slide_timeline(prs, theme, sd, idx, total):
    """Vaqt o'qi (tarix/bosqichlar) — gorizontal chiziq + nuqtalar."""
    slide = _new(prs)
    _scene(slide, theme, idx)
    _title_block(slide, theme, str(sd.get("title", ""))[:110])
    items = [str(p) for p in sd.get("points", []) if str(p).strip()][:4]
    n = max(1, len(items))
    line_y = 4.2
    line = _rect(slide, Inches(0.9), Inches(line_y), Inches(11.5), Inches(0.05))
    _fill(line, theme["accent"], alpha=60)
    seg = 11.5 / n
    for i, it in enumerate(items):
        cx = 0.9 + seg * (i + 0.5)
        dot = _rect(slide, Inches(cx - 0.16), Inches(line_y - 0.13), Inches(0.34), Inches(0.34),
                    MSO_SHAPE.OVAL)
        _fill(dot, theme["accent"])
        above = i % 2 == 0
        card_y = line_y - 1.85 if above else line_y + 0.5
        card = _rect(slide, Inches(cx - seg / 2 + 0.15), Inches(card_y),
                     Inches(seg - 0.3), Inches(1.35), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0] = 0.08
        _fill(card, theme["card"])
        _text(slide, it, Inches(cx - seg / 2 + 0.28), Inches(card_y + 0.08),
             Inches(seg - 0.56), Inches(1.2), size=13, color=theme["text"],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _page_number(slide, theme, idx, total)


def _slide_two_col(prs, theme, sd, idx, total):
    slide = _new(prs)
    _scene(slide, theme, idx)
    _title_block(slide, theme, str(sd.get("title", ""))[:110])
    cols = [c for c in (sd.get("columns") or []) if isinstance(c, dict)][:2]
    col_w = Inches(5.85)
    lefts = [Inches(0.6), Inches(6.85)]
    for ci, col in enumerate(cols):
        left = lefts[ci]
        card = _rect(slide, left, Inches(2.1), col_w, Inches(4.6),
                     shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0] = 0.05
        _fill(card, theme["card"])
        head = _rect(slide, left, Inches(2.1), col_w, Inches(0.8),
                     shape_type=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
        head.adjustments[0] = 0.3
        _fill(head, theme["card2"])
        _text(slide, str(col.get("title", ""))[:60], Inches(left.inches + 0.35), Inches(2.18),
             Inches(col_w.inches - 0.7), Inches(0.65), size=19, bold=True,
             color=theme["accent"], anchor=MSO_ANCHOR.MIDDLE)
        top = 3.15
        for pt in [str(p) for p in col.get("points", [])][:4]:
            h = 0.5 + 0.28 * (_lines_for(pt, 40) - 1)
            _text(slide, f"▪  {pt}", Inches(left.inches + 0.35), Inches(top),
                 Inches(col_w.inches - 0.7), Inches(h), size=15, color=theme["text"])
            top += h + 0.12
    _page_number(slide, theme, idx, total)


def _slide_stat(prs, theme, sd, idx, total):
    slide = _new(prs)
    _scene(slide, theme, idx)
    value = str(sd.get("value", ""))[:20]
    label = str(sd.get("label") or sd.get("title", ""))[:120]
    _text(slide, value, Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.4),
         size=100, bold=True, color=theme["accent"], align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, label, Inches(1.5), Inches(4.6), Inches(10.3), Inches(1.2),
         size=22, color=theme["title"], align=PP_ALIGN.CENTER)
    _page_number(slide, theme, idx, total)


def _slide_quote(prs, theme, sd, idx, total):
    slide = _new(prs)
    _scene(slide, theme, idx)
    _text(slide, "“", Inches(0.8), Inches(0.85), Inches(2.5), Inches(2.0),
         size=110, bold=True, color=theme["accent"])
    quote = str(sd.get("quote_text") or sd.get("title", ""))[:280]
    _text(slide, quote, Inches(1.4), Inches(2.4), Inches(10.5), Inches(2.6),
         size=25, bold=True, color=theme["title"])
    author = str(sd.get("quote_author", ""))[:80]
    if author:
        _text(slide, f"— {author}", Inches(1.4), Inches(5.35), Inches(10.5), Inches(0.6),
             size=17, color=theme["accent"])
    _page_number(slide, theme, idx, total)


def _slide_chart(prs, theme, sd, idx, total, temp_files):
    slide = _new(prs)
    _scene(slide, theme, idx)
    _title_block(slide, theme, str(sd.get("title", ""))[:110])
    spec = sd.get("chart")
    chart_ok = False
    if isinstance(spec, dict):
        try:
            from docs.charts import render_chart
            fd, path = tempfile.mkstemp(suffix=".png")
            os.close(fd)
            render_chart(spec, path)
            temp_files.append(path)
            card = _rect(slide, Inches(0.6), Inches(2.15), Inches(8.3), Inches(4.55),
                        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            card.adjustments[0] = 0.04
            _fill(card, _WHITE)
            slide.shapes.add_picture(path, Inches(0.75), Inches(2.3), width=Inches(8.0))
            chart_ok = True
        except Exception:
            chart_ok = False
    points = [str(p) for p in sd.get("points", [])][:3]
    if chart_ok:
        left, width, cpl, size = Inches(9.15), Inches(3.6), 30, 15
    else:
        left, width, cpl, size = Inches(0.6), Inches(12.1), 68, 17
    top = 2.3
    for pt in points:
        h = 0.55 + 0.3 * (_lines_for(pt, cpl) - 1)
        _text(slide, f"▪ {pt}", left, Inches(top), width, Inches(h), size=size, color=theme["text"])
        top += h + 0.2
    if not chart_ok and not points:
        _text(slide, "Ma'lumot yetarli emas", left, Inches(2.3), width, Inches(0.6),
             size=16, color=theme["text"])
    _page_number(slide, theme, idx, total)


def _slide_image(prs, theme, sd, image_bytes, idx, total):
    slide = _new(prs)
    if image_bytes:
        _cover_photo(slide, image_bytes, overlay_alpha=30)
    else:
        _scene(slide, theme, idx)
    title = str(sd.get("title", ""))[:110]
    band = _rect(slide, 0, Inches(5.6), _SLIDE_W, Inches(1.9))
    _fill(band, _OVERLAY, alpha=72)
    _text(slide, title, Inches(0.7), Inches(5.75), Inches(11.9), Inches(0.8),
         size=26, bold=True, color=_WHITE)
    points = [str(p) for p in sd.get("points", [])][:1]
    if points:
        _text(slide, points[0], Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.9),
             size=16, color=RGBColor(0xe6, 0xe6, 0xeb))
    _page_number(slide, theme, idx, total)


def _slide_end(prs, theme, sd, image_bytes):
    slide = _new(prs)
    if image_bytes:
        _cover_photo(slide, image_bytes, overlay_alpha=62)
    else:
        _bg_gradient(slide, theme, "radial_center")
        _decorate(slide, theme, "blobs")
        _accent_bar(slide, theme, Inches(3.4))
    title = str(sd.get("title", "")) or "E'tiboringiz uchun rahmat!"
    _text(slide, title, Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.6),
         size=40, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
    top = 4.2
    for pt in [str(p) for p in sd.get("points", [])][:3]:
        _text(slide, pt, Inches(2.0), Inches(top), Inches(9.3), Inches(0.5),
             size=17, color=RGBColor(0xe6, 0xe6, 0xeb), align=PP_ALIGN.CENTER)
        top += 0.55


# ─────────────────────── Parsing, rasm olish, yig'ish ───────────────────────

def parse_slides(json_text: str) -> list[dict]:
    data = parse_json_objects(json_text)
    return [d for d in data if d.get("title") or d.get("points") or d.get("chart")]


async def fetch_images_for_slides(slides: list[dict]) -> dict[int, bytes]:
    from docs.pexels import fetch_image

    queries: dict[int, str] = {}
    for i, s in enumerate(slides):
        layout = (s.get("layout") or "").lower()
        is_cover = i == 0 or layout == "cover"
        is_image = layout == "image"
        if is_cover or is_image:
            q = (s.get("image_query") or "").strip() or str(s.get("title", ""))
            if q:
                queries[i] = q
    if not queries:
        return {}

    cache: dict[str, bytes | None] = {}

    async def _get(q: str):
        if q not in cache:
            cache[q] = await fetch_image(q)
        return cache[q]

    results = await asyncio.gather(*[_get(q) for q in queries.values()])
    return {idx: res for (idx, _), res in zip(queries.items(), results) if res}


def build_pptx(slides: list[dict], output_path: str, topic: str = "",
              author: str = "", subject: str = "",
              images: dict[int, bytes] | None = None) -> str:
    if len(slides) < 3:
        raise ValueError(
            f"Slayd ma'lumotlari yetarli emas ({len(slides)} ta) — AI javobi buzuq/kesilgan")
    images = images or {}
    theme = _theme_for(topic)
    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H

    total = len(slides)
    temp_files: list[str] = []
    section_no = 0
    for i, sd in enumerate(slides):
        layout = (sd.get("layout") or "content").lower()
        img = images.get(i)
        n = i + 1
        if i == 0 or layout == "cover":
            _slide_cover(prs, theme, sd, img, author, subject)
        elif i == total - 1 or layout == "end":
            _slide_end(prs, theme, sd, img)
        elif layout == "section":
            section_no += 1
            _slide_section(prs, theme, sd, n, total, section_no)
        elif layout == "steps" and sd.get("points"):
            _slide_steps(prs, theme, sd, n, total)
        elif layout == "cards" and sd.get("points"):
            _slide_cards(prs, theme, sd, n, total)
        elif layout == "timeline" and sd.get("points"):
            _slide_timeline(prs, theme, sd, n, total)
        elif layout == "two_col" and sd.get("columns"):
            _slide_two_col(prs, theme, sd, n, total)
        elif layout == "stat" and (sd.get("value") or sd.get("label")):
            _slide_stat(prs, theme, sd, n, total)
        elif layout == "quote" and sd.get("quote_text"):
            _slide_quote(prs, theme, sd, n, total)
        elif layout == "chart" and sd.get("chart"):
            _slide_chart(prs, theme, sd, n, total, temp_files)
        elif layout == "image":
            _slide_image(prs, theme, sd, img, n, total)
        else:
            _slide_content(prs, theme, sd, n, total)

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    for f in temp_files:
        try:
            os.remove(f)
        except OSError:
            pass
    return output_path
